from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import copy
import os
import uuid

# ============================================================
# Brand Tokens
# ============================================================
BLUE = RGBColor(0x10, 0x32, 0xF5)
BLACK = RGBColor(0x00, 0x00, 0x00)
TITLE_DARK = RGBColor(0x1E, 0x29, 0x3B)
BULLET_GRAY = RGBColor(0x64, 0x74, 0x8B)
SUBTITLE_GRAY = RGBColor(0x94, 0xA3, 0xB8)
CARD_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BADGE_BG = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NUM_GRAY = RGBColor(0x99, 0x99, 0x99)
BEFORE_COLOR = RGBColor(0x64, 0x74, 0x8B)
AFTER_COLOR = BLUE

# ============================================================
# Slide Dimensions & Shared Positions
# ============================================================
SLIDE_W = Emu(24384000)
SLIDE_H = Emu(13716000)

TITLE_POS = (Emu(1177725), Emu(874446), Emu(16000000), Emu(1308101))
PAGE_NUM_POS = (Emu(1474857), Emu(12798000), Emu(400000), Emu(469901))
LOGO_NAME = "Dify-logo.svg"

CARD_AREA_X = Emu(1270000)
CARD_AREA_Y = Emu(3431755)
CARD_AREA_W = Emu(21844000)
CARD_AREA_H = Emu(7785101)
CARD_GAP = Emu(274320)
CARD_PAD = Emu(228600)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(BASE_DIR, "templates", "general-intro.pptx")
OUTPUT_DIR = os.path.join(BASE_DIR, "output")

_ref_prs = None

def _get_ref():
    global _ref_prs
    if _ref_prs is None and os.path.exists(TEMPLATE_PATH):
        _ref_prs = Presentation(TEMPLATE_PATH)
    return _ref_prs


def _add_shadow(shape):
    spPr = shape._element.find('.//' + qn('a:spPr'))
    if spPr is None:
        return
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '38100')
    outerShdw.set('dist', '12700')
    outerShdw.set('dir', '5400000')
    outerShdw.set('algn', 'tl')
    srgb = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgb.set('val', '000000')
    alpha = etree.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', '8000')


def _copy_assets(slide, bg_name, include_logo=True):
    """Copy background image and logo from reference template."""
    ref = _get_ref()
    if ref is None:
        return
    for ref_slide in ref.slides:
        found_bg = False
        found_logo = False
        for shape in ref_slide.shapes:
            if shape.name == bg_name and not found_bg:
                el = copy.deepcopy(shape._element)
                slide.shapes._spTree.insert(2, el)
                for rel in ref_slide.part.rels.values():
                    if hasattr(rel, 'target_part'):
                        try:
                            slide.part.relate_to(rel.target_part, rel.reltype, rel.rId)
                        except:
                            pass
                found_bg = True
            if include_logo and shape.name == LOGO_NAME and not found_logo:
                el = copy.deepcopy(shape._element)
                slide.shapes._spTree.append(el)
                for rel in ref_slide.part.rels.values():
                    if hasattr(rel, 'target_part'):
                        try:
                            slide.part.relate_to(rel.target_part, rel.reltype, rel.rId)
                        except:
                            pass
                found_logo = True
        if found_bg:
            break


def _add_text(slide, x, y, w, h, text, font_name="Söhne Kräftig", size=60,
              color=BLACK, bold=False, alignment=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font_name
    return tf


def _add_title(slide, text_black, text_blue, pos=None, size=60):
    px, py, pw, ph = pos or TITLE_POS
    tb = slide.shapes.add_textbox(px, py, pw, ph)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = text_black + " "
    r1.font.size = Pt(size)
    r1.font.bold = False
    r1.font.color.rgb = BLACK
    r1.font.name = "Söhne Kräftig"
    r2 = p.add_run()
    r2.text = text_blue
    r2.font.size = Pt(size)
    r2.font.bold = False
    r2.font.color.rgb = BLUE
    r2.font.name = "Söhne Kräftig"


def _add_page_number(slide, num):
    px, py, pw, ph = PAGE_NUM_POS
    tb = slide.shapes.add_textbox(px, py, pw, ph)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = str(num)
    r.font.size = Pt(18)
    r.font.bold = False
    r.font.color.rgb = NUM_GRAY
    r.font.name = "Söhne"


def _new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


# ============================================================
# Dynamic font sizing helpers
# ============================================================
def _auto_bullet_size(bullets, font_scale=1.0, base=16, floor=8, ceil=16):
    """Shrink bullet font when there is a lot of text in the card."""
    if not bullets:
        return max(floor, int(base * font_scale))
    total_chars = sum(len(b) for b in bullets)
    n = len(bullets)
    if total_chars <= 120 and n <= 3:
        size = base
    elif total_chars <= 220 and n <= 5:
        size = 14
    elif total_chars <= 350 and n <= 6:
        size = 12
    elif total_chars <= 500:
        size = 11
    elif total_chars <= 700:
        size = 10
    else:
        size = 9
    return max(floor, min(ceil, int(size * font_scale)))


def _auto_cover_title_size(title, base=72, floor=36, ceil=72):
    """Shrink cover title for long text."""
    n = len(title)
    if n <= 40:
        return ceil
    elif n <= 60:
        return 60
    elif n <= 80:
        return 52
    elif n <= 120:
        return 44
    else:
        return max(floor, 36)


def _auto_metric_sizes(metrics, cols):
    """Adaptive number/title/desc sizes for big_numbers based on content length."""
    total_desc = sum(len(m.get("description", "")) for m in metrics)
    if cols <= 2 and total_desc <= 200:
        return 100, 22, 16
    elif cols <= 3 and total_desc <= 400:
        return 80, 20, 14
    else:
        return 64, 18, 12


# ============================================================
# Layout: Cards (2-6 items, auto grid)
# ============================================================
def _draw_card(slide, x, y, w, h, data, font_scale=1.0):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_WHITE
    card.line.fill.background()
    _add_shadow(card)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(36576))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.fill.background()

    badge_s = Emu(int(502920 * font_scale))
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        x + CARD_PAD, y + CARD_PAD, badge_s, badge_s)
    badge.fill.solid()
    badge.fill.fore_color.rgb = BADGE_BG
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = data.get("num", "")
    r.font.size = Pt(int(20 * font_scale))
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Söhne Kräftig"

    title_size = max(16, int(22 * font_scale))
    title_y = y + Emu(int(868680 * font_scale))
    tb = slide.shapes.add_textbox(x + CARD_PAD, title_y, w - CARD_PAD * 2, Emu(640080))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.space_before = Pt(0)
    p1.space_after = Pt(0)
    r1 = p1.add_run()
    r1.text = data.get("title_blue", "")
    r1.font.size = Pt(title_size)
    r1.font.bold = True
    r1.font.color.rgb = BLUE
    r1.font.name = "Söhne Kräftig"
    p2 = tf.add_paragraph()
    p2.space_before = Pt(0)
    p2.space_after = Pt(0)
    r2 = p2.add_run()
    r2.text = data.get("title_black", "")
    r2.font.size = Pt(title_size)
    r2.font.bold = True
    r2.font.color.rgb = TITLE_DARK
    r2.font.name = "Söhne Kräftig"

    bullets = data.get("bullets", [])
    bullet_size = _auto_bullet_size(bullets, font_scale)
    bullet_spacing = max(1, min(4, int(4 * (bullet_size / 16))))
    bullet_y = y + Emu(int(1645920 * font_scale))
    bb = slide.shapes.add_textbox(x + CARD_PAD, bullet_y, w - CARD_PAD * 2, h - (bullet_y - y) - CARD_PAD)
    tf2 = bb.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(bullet_spacing)
        r = p.add_run()
        r.text = "•  " + bullet
        r.font.size = Pt(bullet_size)
        r.font.color.rgb = BULLET_GRAY
        r.font.name = "Söhne"


def build_cards(prs, data):
    slide = _new_slide(prs)
    _copy_assets(slide, "The Challenge Number Background.png")
    _add_title(slide, data.get("title_black", ""), data.get("title_blue", ""))
    _add_page_number(slide, data.get("page_number", 1))

    cards = data.get("cards", [])
    n = len(cards)
    if n == 0:
        return slide

    if n <= 3:
        cols, rows = n, 1
    elif n == 4:
        cols, rows = 4, 1
    elif n <= 6:
        cols, rows = 3, 2
    else:
        cols, rows = 4, 2

    gap = CARD_GAP
    card_w = (CARD_AREA_W - gap * (cols - 1)) // cols
    card_h = (CARD_AREA_H - gap * (rows - 1)) // rows if rows > 1 else CARD_AREA_H

    font_scale = 1.0
    if cols >= 4:
        font_scale = 0.85
    if rows >= 2 and cols >= 3:
        font_scale = 0.9

    for idx, card_data in enumerate(cards[:cols * rows]):
        row = idx // cols
        col = idx % cols

        items_in_row = min(cols, n - row * cols)
        row_w = card_w * items_in_row + gap * (items_in_row - 1)
        row_offset = CARD_AREA_X + (CARD_AREA_W - row_w) // 2

        if row == 0:
            row_offset = CARD_AREA_X
        elif row == 1 and n > cols:
            remaining = n - cols
            row_w = card_w * remaining + gap * (remaining - 1)
            row_offset = CARD_AREA_X + (CARD_AREA_W - row_w) // 2
            col = idx - cols

        x = row_offset + col * (card_w + gap)
        y = CARD_AREA_Y + row * (card_h + gap)
        _draw_card(slide, x, y, card_w, card_h, card_data, font_scale)

    return slide


# ============================================================
# Layout: Big Numbers (2-6 metrics)
# ============================================================
def build_big_numbers(prs, data):
    slide = _new_slide(prs)
    _copy_assets(slide, "The Challenge Number Background.png")
    _add_title(slide, data.get("title_black", ""), data.get("title_blue", ""))
    _add_page_number(slide, data.get("page_number", 1))

    metrics = data.get("metrics", [])
    n = len(metrics)
    if n == 0:
        return slide

    if n <= 2:
        cols, rows = n, 1
    elif n <= 4:
        cols, rows = 2, 2
    elif n <= 6:
        cols, rows = 3, 2
    else:
        cols, rows = 3, 2

    cell_w = CARD_AREA_W // cols
    cell_h = CARD_AREA_H // rows

    number_size, title_size, desc_size = _auto_metric_sizes(metrics, cols)
    unit_size = int(number_size * 0.45)

    for idx, m in enumerate(metrics[:cols * rows]):
        row = idx // cols
        col = idx % cols
        cx = CARD_AREA_X + col * cell_w
        cy = CARD_AREA_Y + row * cell_h

        left_w = Emu(int(cell_w * 0.4))
        right_x = cx + left_w
        right_w = cell_w - left_w

        num_text = m.get("number", "0")
        unit_text = m.get("unit", "")

        tb = slide.shapes.add_textbox(cx, cy + Emu(200000), left_w, Emu(int(cell_h * 0.6)))
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = num_text
        r.font.size = Pt(number_size)
        r.font.bold = True
        r.font.color.rgb = BLACK
        r.font.name = "Söhne Kräftig"

        if unit_text:
            r2 = p.add_run()
            r2.text = unit_text
            r2.font.size = Pt(unit_size)
            r2.font.bold = False
            r2.font.color.rgb = BLACK
            r2.font.name = "Söhne Kräftig"

        _add_text(slide, right_x, cy + Emu(150000), right_w - CARD_PAD, Emu(600000),
                  m.get("title", ""), size=title_size, bold=True, color=TITLE_DARK)

        _add_text(slide, right_x, cy + Emu(750000), right_w - CARD_PAD, Emu(int(cell_h * 0.5)),
                  m.get("description", ""), font_name="Söhne", size=desc_size, color=BULLET_GRAY)

    return slide


# ============================================================
# Layout: Cover / Hero
# ============================================================
def build_cover(prs, data):
    slide = _new_slide(prs)

    ref = _get_ref()
    if ref and len(ref.slides) >= 3:
        ref_slide = ref.slides[2]
        for shape in ref_slide.shapes:
            if "背景" in shape.name or "background" in shape.name.lower():
                el = copy.deepcopy(shape._element)
                slide.shapes._spTree.insert(2, el)
                for rel in ref_slide.part.rels.values():
                    if hasattr(rel, 'target_part'):
                        try:
                            slide.part.relate_to(rel.target_part, rel.reltype, rel.rId)
                        except:
                            pass
                break
        for shape in ref_slide.shapes:
            if "logo" in shape.name.lower() and "dify" in shape.name.lower():
                el = copy.deepcopy(shape._element)
                slide.shapes._spTree.append(el)
                for rel in ref_slide.part.rels.values():
                    if hasattr(rel, 'target_part'):
                        try:
                            slide.part.relate_to(rel.target_part, rel.reltype, rel.rId)
                        except:
                            pass
                break

    title = data.get("title", "")
    title_blue = data.get("title_blue", "")
    cover_size = _auto_cover_title_size(title)

    tb = slide.shapes.add_textbox(Emu(1222109), Emu(491817), Emu(13747191), Emu(3188971))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    if title_blue and title_blue in title:
        parts = title.split(title_blue, 1)
        if parts[0]:
            r = p.add_run()
            r.text = parts[0]
            r.font.size = Pt(cover_size)
            r.font.bold = False
            r.font.color.rgb = BLACK
            r.font.name = "Söhne Kräftig"
        r_blue = p.add_run()
        r_blue.text = title_blue
        r_blue.font.size = Pt(cover_size)
        r_blue.font.bold = False
        r_blue.font.color.rgb = BLUE
        r_blue.font.name = "Söhne Kräftig"
        if len(parts) > 1 and parts[1]:
            r_after = p.add_run()
            r_after.text = parts[1]
            r_after.font.size = Pt(cover_size)
            r_after.font.bold = False
            r_after.font.color.rgb = BLACK
            r_after.font.name = "Söhne Kräftig"
    else:
        r = p.add_run()
        r.text = title
        r.font.size = Pt(cover_size)
        r.font.bold = False
        r.font.color.rgb = BLACK
        r.font.name = "Söhne Kräftig"

    subtitle = data.get("subtitle", "")
    if subtitle:
        sub_size = 24 if len(subtitle) <= 80 else 20 if len(subtitle) <= 140 else 16
        _add_text(slide, Emu(1221068), Emu(3534821), Emu(11341057), Emu(1187451),
                  subtitle, font_name="Söhne", size=sub_size, color=BULLET_GRAY)

    footnote = data.get("footnote", "")
    if footnote:
        fn_size = 14 if len(footnote) <= 120 else 12 if len(footnote) <= 250 else 10
        _add_text(slide, Emu(1245121), Emu(12375667), Emu(18934803), Emu(978866),
                  footnote, font_name="Söhne", size=fn_size, color=SUBTITLE_GRAY)

    return slide


# ============================================================
# Layout: Case Study (Before / After)
# ============================================================
def build_case_study(prs, data):
    slide = _new_slide(prs)
    _copy_assets(slide, "5 (2).png")

    title_black = data.get("title_black", "From Business Challenge")
    title_blue = data.get("title_blue", "to Tangible Results")
    _add_title(slide, title_black, title_blue)

    subtitle = data.get("subtitle", "")
    if subtitle:
        _add_text(slide, Emu(1314450), Emu(2889250), Emu(5013046), Emu(749301),
                  subtitle, font_name="Söhne", size=20, color=BULLET_GRAY)

    company_desc = data.get("company_description", "")
    if company_desc:
        _add_text(slide, Emu(1379157), Emu(4800000), Emu(4023180), Emu(2500000),
                  company_desc, font_name="Söhne", size=18, color=SUBTITLE_GRAY)

    header_y = Emu(4234765)
    before_x = Emu(6879696)
    after_x = Emu(13873894)
    before_w = Emu(5461001)
    after_w = Emu(9105901)
    header_h = Emu(571501)

    before_hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, before_x, header_y, before_w, header_h)
    before_hdr.fill.solid()
    before_hdr.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
    before_hdr.line.fill.background()
    top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, before_x, header_y, before_w, Emu(36576))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = SUBTITLE_GRAY
    top_line.line.fill.background()
    _add_text(slide, before_x + CARD_PAD, header_y + Emu(100000), before_w, Emu(400000),
              "BEFORE  ·  CHALLENGE", font_name="Söhne Kräftig", size=14, color=SUBTITLE_GRAY)

    after_hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, after_x, header_y, after_w, header_h)
    after_hdr.fill.solid()
    after_hdr.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xFF)
    after_hdr.line.fill.background()
    top_line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, after_x, header_y, after_w, Emu(36576))
    top_line2.fill.solid()
    top_line2.fill.fore_color.rgb = BLUE
    top_line2.line.fill.background()
    _add_text(slide, after_x + CARD_PAD, header_y + Emu(100000), after_w, Emu(400000),
              "AFTER  ·  SOLUTION", font_name="Söhne Kräftig", size=14, color=BLUE)

    rows = data.get("rows", [])
    n = len(rows)
    if n == 0:
        return slide

    row_start_y = Emu(5000000)
    available_h = Emu(13716000 - 5000000 - 1500000)
    row_h = available_h // max(n, 1)
    row_gap = Emu(100000)
    arrow_x = Emu(12500000)

    for idx, row in enumerate(rows):
        ry = row_start_y + idx * (row_h + row_gap)

        before_title = row.get("before_title", "")
        before_desc = row.get("before_description", "")
        after_title = row.get("after_title", "")
        after_desc = row.get("after_description", "")
        tag = row.get("tag", "")

        _add_text(slide, before_x + CARD_PAD, ry, before_w - CARD_PAD * 2, Emu(350000),
                  before_title, font_name="Söhne Kräftig", size=18, bold=True, color=TITLE_DARK)
        _add_text(slide, before_x + CARD_PAD, ry + Emu(400000), before_w - CARD_PAD * 2, Emu(600000),
                  before_desc, font_name="Söhne", size=15, color=BULLET_GRAY)

        _add_text(slide, arrow_x, ry + Emu(150000), Emu(800000), Emu(300000),
                  "→", font_name="Söhne", size=24, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)

        _add_text(slide, after_x + CARD_PAD, ry, after_w - CARD_PAD * 4, Emu(350000),
                  after_title, font_name="Söhne Kräftig", size=18, bold=True, color=TITLE_DARK)

        if tag:
            tag_x = after_x + after_w - Emu(3200000)
            tag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tag_x, ry + Emu(30000), Emu(3000000), Emu(300000))
            tag_box.fill.solid()
            tag_box.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
            tag_box.line.fill.background()
            _add_text(slide, tag_x + Emu(50000), ry + Emu(60000), Emu(2900000), Emu(250000),
                      tag, font_name="Söhne Kräftig", size=11, color=TITLE_DARK, alignment=PP_ALIGN.CENTER)

        _add_text(slide, after_x + CARD_PAD, ry + Emu(400000), after_w - CARD_PAD * 2, Emu(600000),
                  after_desc, font_name="Söhne", size=15, color=BULLET_GRAY)

    _add_page_number(slide, data.get("page_number", 1))
    return slide


# ============================================================
# Router: build_from_json
# ============================================================
def build_from_json(data: dict) -> str:
    """
    Build a PPT slide from structured JSON.

    The "layout_type" field determines which layout to use:
    - "cards"       : Card grid (2-6 cards, auto rows/cols)
    - "big_numbers" : Large metrics display (2-6 stats)
    - "cover"       : Hero/cover slide
    - "case_study"  : Before/after case study

    Returns: path to the generated .pptx file.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    layout_type = data.get("layout_type", "cards")

    if layout_type == "cards":
        build_cards(prs, data)
    elif layout_type == "big_numbers":
        build_big_numbers(prs, data)
    elif layout_type == "cover":
        build_cover(prs, data)
    elif layout_type == "case_study":
        build_case_study(prs, data)
    else:
        build_cards(prs, data)

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    path = os.path.join(OUTPUT_DIR, f"slide-{uuid.uuid4().hex[:8]}.pptx")
    prs.save(path)
    return path
