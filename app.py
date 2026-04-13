import os
import json
import uuid
import base64
import subprocess
import logging
from datetime import datetime
from flask import Flask, request, jsonify, send_file, render_template_string
from anthropic import Anthropic
from builder import build_from_json
from pptx import Presentation as ReadPptx
import memory_store

app = Flask(__name__)
client = Anthropic()
memory_store.init_db()

LOG_DIR = os.path.join(os.path.dirname(__file__), 'logs')
os.makedirs(LOG_DIR, exist_ok=True)
file_handler = logging.FileHandler(os.path.join(LOG_DIR, 'app.log'))
file_handler.setFormatter(logging.Formatter('%(asctime)s [%(levelname)s] %(message)s'))
app.logger.addHandler(file_handler)
app.logger.setLevel(logging.INFO)

sessions = {}

SYSTEM_PROMPT = """You are a Dify PPT slide generator. You output ONLY valid JSON (no markdown, no explanation).

First, choose the best layout_type based on the user's content:

## Layout Types & Schemas

### 1. "cards" — Feature cards (2-6 items with titles and bullets)
Best for: product capabilities, solution pillars, feature lists
```
{
  "layout_type": "cards",
  "title_black": "What we do",
  "title_blue": "Today",
  "page_number": 1,
  "cards": [
    {"num": "1", "title_blue": "End-to-End", "title_black": "AI Platform", "bullets": ["bullet 1", "bullet 2"]}
  ]
}
```
Grid auto-adjusts: 2-3 items = 1 row, 4 items = 4 cols, 5-6 items = 3+2 or 3+3.

### 2. "big_numbers" — Key metrics / statistics (2-6 numbers)
Best for: KPIs, growth stats, impact metrics, market data
```
{
  "layout_type": "big_numbers",
  "title_black": "Community-Driven",
  "title_blue": "Globally Recognized",
  "page_number": 1,
  "metrics": [
    {"number": "1M+", "unit": "", "title": "Apps Powered by Dify", "description": "Over one million applications built on the platform"}
  ]
}
```
"unit" is optional (use for %, x, etc. shown next to the number).

### 3. "cover" — Hero / title slide
Best for: first slide, section openers, keynote titles
```
{
  "layout_type": "cover",
  "title": "Build Production-Ready Agentic Workflows",
  "title_blue": "Agentic Workflows",
  "subtitle": "Develop, deploy, and manage agentic workflows for teams at any scale.",
  "footnote": "Prepared for [audience] — March 2026",
  "page_number": 1
}
```
"title_blue" is the keyword within "title" that should be highlighted in blue.

### 4. "case_study" — Before/After comparison
Best for: customer success stories, transformation results
```
{
  "layout_type": "case_study",
  "title_black": "From Business Challenge",
  "title_blue": "to Tangible Results",
  "subtitle": "Case Study on CompanyName",
  "company_description": "Brief description of the company...",
  "page_number": 1,
  "rows": [
    {
      "before_title": "Limited AI Expertise",
      "before_description": "Description of the challenge...",
      "after_title": "Self-Service AI",
      "after_description": "Description of the solution...",
      "tag": "DEMOCRATIZED ACCESS"
    }
  ]
}
```
"tag" is an optional short label shown as a badge on the after side. 2-4 rows recommended.

## Rules
- Choose the layout_type that best matches the user's content
- ALL text must be in English
- Keep text concise — titles under 40 chars, bullets under 60 chars
- If modifying a previous slide, output the COMPLETE updated JSON
- If the user provides an image, extract content and choose the best layout
- If the user uploads a file, use the extracted text to build the slide
- If unclear, default to "cards" layout
- Output exactly ONE JSON object (root must be `{...}`), never a JSON array `[...]` of slides — this tool generates one slide per request."""


def normalize_slide_json(parsed):
    """
    Model sometimes returns [{slide1},{slide2}] or {slides:[...]}.
    This app is single-slide only: pick the first valid slide dict.
    Returns (slide_dict, user_note) — user_note explains if extra slides were dropped.
    """
    dropped = ""
    if isinstance(parsed, list):
        if not parsed:
            raise ValueError("JSON array is empty; need one slide object.")
        if len(parsed) > 1:
            dropped = "\n\nNote: The model returned multiple slides; only the first slide was exported (this tool is single-slide per request)."
        for item in parsed:
            if isinstance(item, dict) and item.get("layout_type"):
                return item, dropped
        first = parsed[0]
        if isinstance(first, dict):
            return first, dropped
        raise ValueError("JSON array must contain slide objects.")
    if isinstance(parsed, dict):
        slides = parsed.get("slides")
        if isinstance(slides, list) and slides:
            if len(slides) > 1:
                dropped = "\n\nNote: Multiple slides were in the response; only the first was exported (single-slide mode)."
            for item in slides:
                if isinstance(item, dict) and item.get("layout_type"):
                    return item, dropped
            if isinstance(slides[0], dict):
                return slides[0], dropped
        return parsed, ""
    raise ValueError("JSON root must be an object or array of objects.")


def build_system_prompt():
    """Base rules + compact memory block (set MEMORY_DISABLE=1 to skip)."""
    if os.environ.get("MEMORY_DISABLE", "").lower() in ("1", "true", "yes"):
        return SYSTEM_PROMPT
    return SYSTEM_PROMPT + "\n\n" + memory_store.format_memory_block()


def _anthropic_model_chain():
    """Primary + fallback(s). ANTHROPIC_MODEL_FALLBACK can be comma-separated."""
    primary = os.environ.get("ANTHROPIC_MODEL", "claude-opus-4-5").strip()
    raw_fb = os.environ.get(
        "ANTHROPIC_MODEL_FALLBACK",
        "claude-3-5-sonnet-20241022",
    ).strip()
    fallbacks = [x.strip() for x in raw_fb.split(",") if x.strip()]
    seen = set()
    chain = []
    for m in [primary] + fallbacks:
        if m and m not in seen:
            seen.add(m)
            chain.append(m)
    return chain


def messages_create_with_fallback(**api_kwargs):
    """Try ANTHROPIC_MODEL, then each ANTHROPIC_MODEL_FALLBACK entry on API errors."""
    last_err = None
    for i, model_id in enumerate(_anthropic_model_chain()):
        try:
            if i > 0:
                app.logger.warning("anthropic: switching to backup model %s", model_id)
            return client.messages.create(model=model_id, **api_kwargs)
        except Exception as e:
            last_err = e
            app.logger.warning("anthropic: model %s failed: %s", model_id, e)
    raise last_err


def extract_pptx_text(filepath):
    prs = ReadPptx(filepath)
    lines = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t:
                        slide_texts.append(t)
        if slide_texts:
            lines.append(f"[Slide {i+1}]")
            lines.extend(slide_texts)
    return "\n".join(lines)


def extract_pdf_text(filepath):
    try:
        result = subprocess.run(
            ["python3", "-c", f"""
import sys
try:
    import pdfplumber
    with pdfplumber.open("{filepath}") as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text:
                print(f"[Page {{i+1}}]")
                print(text)
except ImportError:
    print("[PDF text extraction requires pdfplumber. Install with: pip3 install pdfplumber]")
"""],
            capture_output=True, text=True, timeout=30
        )
        return result.stdout or "[Could not extract text from PDF. Try uploading as image/screenshot instead.]"
    except Exception:
        return "[Could not extract text from PDF. Try uploading as image/screenshot instead.]"


def extract_text_file(filepath):
    with open(filepath, 'r', errors='ignore') as f:
        return f.read()[:10000]

HTML_TEMPLATE = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Dify PPT Generator</title>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600&display=swap" rel="stylesheet">
<style>
* { margin: 0; padding: 0; box-sizing: border-box; }
body { font-family: 'Inter', system-ui, sans-serif; background: #F2F4F7; height: 100vh; display: flex; flex-direction: column; }
.header { background: white; padding: 16px 24px; border-bottom: 1px solid #E2E8F0; display: flex; align-items: center; gap: 12px; }
.header svg { height: 24px; }
.header h1 { font-size: 18px; font-weight: 600; color: #1E293B; }
.chat-area { flex: 1; overflow-y: auto; padding: 24px; display: flex; flex-direction: column; gap: 16px; }
.msg { max-width: 80%; padding: 12px 16px; border-radius: 12px; font-size: 14px; line-height: 1.6; word-wrap: break-word; }
.msg.user { background: #1032F5; color: white; align-self: flex-end; border-bottom-right-radius: 4px; }
.msg.bot { background: white; color: #1E293B; align-self: flex-start; border-bottom-left-radius: 4px; box-shadow: 0 1px 3px rgba(0,0,0,0.08); }
.msg.bot a { color: #1032F5; font-weight: 500; text-decoration: none; }
.msg.bot a:hover { text-decoration: underline; }
.msg img.preview { max-width: 200px; border-radius: 8px; margin-top: 8px; }
.input-area { background: white; padding: 16px 24px; border-top: 1px solid #E2E8F0; display: flex; gap: 12px; align-items: center; }
.input-area input[type=text] { flex: 1; padding: 12px 16px; border: 1px solid #E2E8F0; border-radius: 8px; font-size: 14px; font-family: inherit; outline: none; }
.input-area input[type=text]:focus { border-color: #1032F5; }
.input-area button { padding: 12px 24px; background: #1032F5; color: white; border: none; border-radius: 8px; font-size: 14px; font-weight: 500; cursor: pointer; font-family: inherit; white-space: nowrap; }
.input-area button:hover { background: #0c28c4; }
.input-area button:disabled { background: #94A3B8; cursor: not-allowed; }
.upload-btn { padding: 10px 14px; background: #F2F4F7; border: 1px solid #E2E8F0; border-radius: 8px; cursor: pointer; font-size: 18px; line-height: 1; }
.upload-btn:hover { background: #E2E8F0; }
.typing { color: #94A3B8; font-style: italic; }
.img-badge { display: inline-block; background: #E8F0FE; color: #1032F5; padding: 4px 10px; border-radius: 6px; font-size: 12px; font-weight: 500; margin-bottom: 6px; }
#imgPreview { display: none; align-items: center; gap: 8px; padding: 8px 24px; background: #FAFBFF; border-top: 1px solid #E2E8F0; }
#imgPreview img { height: 48px; border-radius: 6px; }
#imgPreview .remove { cursor: pointer; color: #94A3B8; font-size: 18px; }
#imgPreview .remove:hover { color: #EF4444; }
</style>
</head>
<body>
<div class="header">
  <svg viewBox="0 0 200 89" fill="none" xmlns="http://www.w3.org/2000/svg">
    <path d="M88.3295 15.5795C94.0043 15.5795 96.1098 12.0962 96.1098 7.78976C96.1098 3.48331 94.0143 0 88.3295 0C82.6448 0 80.5493 3.48331 80.5493 7.78976C80.5493 12.0962 82.6448 15.5795 88.3295 15.5795Z" fill="#1032F5"/>
    <path d="M23.6114 0.0100098H0V72.3162H23.6114C52.7772 72.3162 61.1089 55.6325 61.1089 36.1581C61.1089 16.6837 52.7772 0.0100098 23.6114 0.0100098ZM23.8921 61.1938H13.3347V11.1325H23.8921C40.6657 11.1325 47.7842 19.374 47.7842 36.1581C47.7842 52.9422 40.6657 61.1837 23.8921 61.1837V61.1938Z" fill="#1E293B"/>
    <path d="M115.55 17.2559V22.265H102.777V33.3875H115.55V61.2038H94.4354V22.255H66.6632V33.3775H82.2137V61.1938H63.876V72.3162H147.203V61.1938H127.762V33.3775H147.203V22.255H127.762V11.1325H147.203V0.0100098H132.755C123.26 0.0100098 115.53 7.74958 115.53 17.2559H115.55Z" fill="#1032F5"/>
    <path d="M186.806 22.255L175.697 58.4131L164.588 22.255H151.394L167.505 69.0237C169.18 73.8923 166.342 77.8775 161.199 77.8775H155.554V89H163.856C171.095 89 177.612 84.4125 180.058 77.5964L200 22.255H186.806Z" fill="#1E293B"/>
  </svg>
  <h1>PPT Generator</h1>
</div>
<div class="chat-area" id="chat">
  <div class="msg bot">
    Hi! I'm the Dify PPT Generator.<br><br>
    <b>Single slide only</b> &mdash; each request creates <b>one</b> .pptx with one slide. For a full deck, generate one slide at a time (or ask for the slide you care about first).<br><br>
    Describe the slide you'd like to create and I'll generate a brand-compliant PowerPoint for you.<br><br>
    <b>What I need from you:</b><br>
    &bull; The content for your slide &mdash; title, key points, sections<br>
    &bull; You can type text, paste/upload an image, or upload a file<br><br>
    <b>Supported inputs:</b><br>
    &bull; <b>Text</b> &mdash; describe the slide content<br>
    &bull; <b>Image</b> &mdash; screenshot or photo of an existing slide (I'll read and restructure it)<br>
    &bull; <b>File</b> &mdash; .pptx, .pdf, .txt, .md (I'll extract the content and build a slide from it)<br><br>
    <b>Available layouts:</b><br>
    &bull; <b>Cards</b> &mdash; 2-6 feature cards with titles and bullets<br>
    &bull; <b>Big Numbers</b> &mdash; key metrics and statistics<br>
    &bull; <b>Cover</b> &mdash; hero/title slide<br>
    &bull; <b>Case Study</b> &mdash; before/after comparison<br><br>
    <i>Example: "Create a slide titled 'What we do Today' with 5 product modules covering AI platform, agents, enterprise infra, observability, and multimodal data."</i>
  </div>
</div>
<div id="batchBar" style="padding:10px 24px;background:#FAFBFF;border-top:1px solid #E2E8F0;display:flex;align-items:center;gap:12px;">
  <label style="padding:8px 16px;background:#1032F5;color:white;border-radius:8px;font-size:13px;font-weight:500;cursor:pointer;white-space:nowrap;">
    Batch: Upload .pptx
    <input type="file" id="batchInput" accept=".pptx" style="display:none" onchange="batchUpload(this)">
  </label>
  <span id="batchStatus" style="font-size:13px;color:#64748B;"></span>
</div>
<div id="imgPreview">
  <img id="previewImg" src="">
  <span id="previewLabel" style="font-size:12px;color:#64748B;">Image attached</span>
  <span class="remove" onclick="clearAttachment()">&times;</span>
</div>
<div class="input-area">
  <label class="upload-btn" title="Upload image or file (.pptx, .pdf, .txt, .md)">
    <input type="file" id="fileInput" accept="image/*,.pptx,.pdf,.txt,.md,.csv" style="display:none" onchange="handleFile(this)">
    &#128206;
  </label>
  <input type="text" id="input" placeholder="One slide per request — describe this slide, or paste an image..." autofocus>
  <button id="send" onclick="sendMsg()">Send</button>
</div>
<script>
const chat = document.getElementById('chat');
const input = document.getElementById('input');
const sendBtn = document.getElementById('send');
let sessionId = null;
let pendingImage = null;
let pendingFile = null;

input.addEventListener('keydown', e => { if (e.key === 'Enter' && !sendBtn.disabled) sendMsg(); });

input.addEventListener('paste', e => {
  const items = e.clipboardData.items;
  for (let i = 0; i < items.length; i++) {
    if (items[i].type.startsWith('image/')) {
      const file = items[i].getAsFile();
      processImageFile(file);
      e.preventDefault();
      break;
    }
  }
});

function handleFile(el) {
  if (!el.files[0]) return;
  const file = el.files[0];
  if (file.type.startsWith('image/')) {
    processImageFile(file);
  } else {
    processDocFile(file);
  }
  el.value = '';
}

function processImageFile(file) {
  pendingFile = null;
  const reader = new FileReader();
  reader.onload = e => {
    pendingImage = { data: e.target.result.split(',')[1], type: file.type };
    document.getElementById('previewImg').src = e.target.result;
    document.getElementById('previewImg').style.display = 'block';
    document.getElementById('previewLabel').textContent = 'Image attached: ' + file.name;
    document.getElementById('imgPreview').style.display = 'flex';
  };
  reader.readAsDataURL(file);
}

function processDocFile(file) {
  pendingImage = null;
  pendingFile = file;
  document.getElementById('previewImg').style.display = 'none';
  document.getElementById('previewLabel').textContent = 'File attached: ' + file.name;
  document.getElementById('imgPreview').style.display = 'flex';
}

function clearAttachment() {
  pendingImage = null;
  pendingFile = null;
  document.getElementById('imgPreview').style.display = 'none';
}

async function sendMsg() {
  const text = input.value.trim();
  if (!text && !pendingImage && !pendingFile) return;
  input.value = '';
  sendBtn.disabled = true;

  let userHtml = '';
  if (pendingImage) userHtml += '<span class="img-badge">Image attached</span><br>';
  if (pendingFile) userHtml += '<span class="img-badge">File: ' + esc(pendingFile.name) + '</span><br>';
  if (text) userHtml += esc(text);
  chat.innerHTML += '<div class="msg user">' + userHtml + '</div>';
  chat.innerHTML += '<div class="msg bot typing" id="typing">Generating...</div>';
  chat.scrollTop = chat.scrollHeight;

  let res;
  try {
    if (pendingFile) {
      const formData = new FormData();
      formData.append('message', text || 'Extract the content from this file and create a Dify-branded slide.');
      formData.append('session_id', sessionId || '');
      formData.append('file', pendingFile);
      clearAttachment();
      res = await fetch('/api/chat-file', { method: 'POST', body: formData });
    } else {
      const payload = { message: text || 'Please redesign this slide following Dify brand guidelines.', session_id: sessionId };
      if (pendingImage) payload.image = pendingImage;
      clearAttachment();
      res = await fetch('/api/chat', {
        method: 'POST',
        headers: {'Content-Type': 'application/json'},
        body: JSON.stringify(payload)
      });
    }
    const data = await res.json();
    sessionId = data.session_id;
    document.getElementById('typing').remove();

    let html = esc(data.reply);
    if (data.file_url) {
      html += '<br><br><a href="' + data.file_url + '" download>Download PPT</a>';
    }
    chat.innerHTML += '<div class="msg bot">' + html + '</div>';
  } catch(e) {
    document.getElementById('typing').remove();
    chat.innerHTML += '<div class="msg bot">Error: ' + esc(e.message) + '</div>';
  }
  sendBtn.disabled = false;
  chat.scrollTop = chat.scrollHeight;
  input.focus();
}

function esc(s) { const d = document.createElement('div'); d.textContent = s; return d.innerHTML.replace(/\\n/g, '<br>'); }

async function batchUpload(el) {
  if (!el.files[0]) return;
  const file = el.files[0];
  const status = document.getElementById('batchStatus');
  status.textContent = 'Processing ' + file.name + '...';
  el.disabled = true;
  const fd = new FormData();
  fd.append('file', file);
  try {
    const res = await fetch('/api/batch', { method: 'POST', body: fd });
    const data = await res.json();
    const msg = document.createElement('div');
    msg.className = 'msg bot';
    let html = esc(data.reply || 'Done');
    if (data.file_url) html += '<br><br><a href="' + data.file_url + '">Download branded PPTX</a>';
    msg.innerHTML = html;
    chat.appendChild(msg);
    chat.scrollTop = chat.scrollHeight;
    status.textContent = 'Done! ' + (data.reply || '');
  } catch(e) {
    status.textContent = 'Error: ' + e.message;
  }
  el.disabled = false;
  el.value = '';
}
</script>
</body>
</html>"""


@app.route('/')
def index():
    return render_template_string(HTML_TEMPLATE)


@app.route('/api/chat', methods=['POST'])
def chat_api():
    session_id = None
    try:
        data = request.json
        message = data.get('message', '')
        image_data = data.get('image')
        session_id = data.get('session_id') or str(uuid.uuid4())
        has_image = "yes" if image_data else "no"
        app.logger.info(f"[chat] session={session_id} image={has_image} message={message[:200]}")

        if session_id not in sessions:
            sessions[session_id] = []

        history = sessions[session_id]

        if image_data:
            content = [
                {"type": "image", "source": {"type": "base64", "media_type": image_data["type"], "data": image_data["data"]}},
                {"type": "text", "text": message},
            ]
        else:
            content = message

        history.append({"role": "user", "content": content})

        response = messages_create_with_fallback(
            max_tokens=4096,
            system=build_system_prompt(),
            messages=history,
        )

        assistant_text = response.content[0].text.strip()
        history.append({"role": "assistant", "content": assistant_text})

        try:
            cleaned = assistant_text
            if cleaned.startswith("```"):
                cleaned = cleaned.split("\n", 1)[1]
                if cleaned.endswith("```"):
                    cleaned = cleaned[:-3]
            json_data, multi_note = normalize_slide_json(json.loads(cleaned))

            pptx_path = build_from_json(json_data)
            filename = os.path.basename(pptx_path)

            layout_type = json_data.get("layout_type", "cards")
            summary_parts = [f"Layout: {layout_type}"]

            if layout_type == "cards":
                for c in json_data.get("cards", []):
                    summary_parts.append(f"  {c.get('num','')}. {c.get('title_blue','')} {c.get('title_black','')}")
            elif layout_type == "big_numbers":
                for m in json_data.get("metrics", []):
                    summary_parts.append(f"  {m.get('number','')} — {m.get('title','')}")
            elif layout_type == "cover":
                summary_parts.append(f"  Title: {json_data.get('title','')}")
                if json_data.get('subtitle'):
                    summary_parts.append(f"  Subtitle: {json_data['subtitle']}")
            elif layout_type == "case_study":
                for r in json_data.get("rows", []):
                    summary_parts.append(f"  {r.get('before_title','')} → {r.get('after_title','')}")

            reply = (
                f"PPT generated!\n\n"
                + "\n".join(summary_parts)
                + multi_note
                + f"\n\nTo modify, just tell me what to change (e.g. \"change the title\", \"add a metric\", \"switch to cover layout\")."
            )

            app.logger.info(f"[chat] OK session={session_id} layout={layout_type} file={filename}")
            memory_store.record_event(
                session_id=session_id,
                endpoint="chat",
                has_image=bool(image_data),
                message_preview=message,
                layout_type=layout_type,
                success=True,
                output_file=filename,
            )
            memory_store.maybe_refresh_insights_llm(client)
            return jsonify({
                "reply": reply,
                "file_url": f"/download/{filename}",
                "session_id": session_id,
            })

        except Exception as e:
            app.logger.error(f"[chat] PARSE_ERR session={session_id} error={e}\nmodel_output={assistant_text[:500]}")
            memory_store.record_event(
                session_id=session_id,
                endpoint="chat",
                has_image=bool(image_data),
                message_preview=message,
                success=False,
                error_type="parse_or_build",
                error_detail=str(e),
            )
            return jsonify({
                "reply": f"Generation error: {e}\n\nModel returned: {assistant_text}",
                "file_url": None,
                "session_id": session_id,
            })

    except Exception as e:
        app.logger.error(f"[chat] API_ERR session={session_id} error={e}", exc_info=True)
        memory_store.record_event(
            session_id=session_id,
            endpoint="chat",
            success=False,
            error_type="api",
            error_detail=str(e),
        )
        return jsonify({
            "reply": f"Error: {e}",
            "file_url": None,
            "session_id": session_id or str(uuid.uuid4()),
        })


@app.route('/api/chat-file', methods=['POST'])
def chat_file_api():
    session_id = None
    try:
        message = request.form.get('message', '')
        session_id = request.form.get('session_id') or str(uuid.uuid4())
        uploaded = request.files.get('file')
        app.logger.info(f"[chat-file] session={session_id} file={uploaded.filename if uploaded else 'none'} message={message[:200]}")

        if session_id not in sessions:
            sessions[session_id] = []

        extracted_text = ""
        if uploaded:
            upload_dir = os.path.join(os.path.dirname(__file__), 'output', 'uploads')
            os.makedirs(upload_dir, exist_ok=True)
            filename = f"{uuid.uuid4().hex[:8]}_{uploaded.filename}"
            filepath = os.path.join(upload_dir, filename)
            uploaded.save(filepath)

            ext = os.path.splitext(uploaded.filename)[1].lower()
            if ext == '.pptx':
                extracted_text = extract_pptx_text(filepath)
            elif ext == '.pdf':
                extracted_text = extract_pdf_text(filepath)
            elif ext in ('.txt', '.md', '.csv'):
                extracted_text = extract_text_file(filepath)
            else:
                extracted_text = f"[Unsupported file type: {ext}. Supported: .pptx, .pdf, .txt, .md]"

        full_message = message
        if extracted_text:
            full_message += f"\n\n--- Extracted file content ---\n{extracted_text[:8000]}\n--- End of file content ---"

        history = sessions[session_id]
        history.append({"role": "user", "content": full_message})

        response = messages_create_with_fallback(
            max_tokens=4096,
            system=build_system_prompt(),
            messages=history,
        )

        assistant_text = response.content[0].text.strip()
        history.append({"role": "assistant", "content": assistant_text})

        try:
            cleaned = assistant_text
            if cleaned.startswith("```"):
                cleaned = cleaned.split("\n", 1)[1]
                if cleaned.endswith("```"):
                    cleaned = cleaned[:-3]
            json_data, multi_note = normalize_slide_json(json.loads(cleaned))

            pptx_path = build_from_json(json_data)
            out_filename = os.path.basename(pptx_path)

            layout_type = json_data.get("layout_type", "cards")

            reply = (
                f"PPT generated from your file!\n\n"
                f"Layout: {layout_type}"
                + multi_note
                + "\n\nTo modify, just tell me what to change."
            )

            app.logger.info(f"[chat-file] OK session={session_id} layout={layout_type} file={out_filename}")
            memory_store.record_event(
                session_id=session_id,
                endpoint="chat-file",
                has_image=False,
                message_preview=message,
                layout_type=layout_type,
                success=True,
                output_file=out_filename,
            )
            memory_store.maybe_refresh_insights_llm(client)
            return jsonify({
                "reply": reply,
                "file_url": f"/download/{out_filename}",
                "session_id": session_id,
            })

        except Exception as e:
            app.logger.error(f"[chat-file] PARSE_ERR session={session_id} error={e}\nmodel_output={assistant_text[:500]}")
            memory_store.record_event(
                session_id=session_id,
                endpoint="chat-file",
                has_image=False,
                message_preview=message,
                success=False,
                error_type="parse_or_build",
                error_detail=str(e),
            )
            return jsonify({
                "reply": f"Generation error: {e}\n\nModel returned: {assistant_text}",
                "file_url": None,
                "session_id": session_id,
            })

    except Exception as e:
        app.logger.error(f"[chat-file] API_ERR session={session_id} error={e}", exc_info=True)
        memory_store.record_event(
            session_id=session_id,
            endpoint="chat-file",
            success=False,
            error_type="api",
            error_detail=str(e),
        )
        return jsonify({
            "reply": f"Error: {e}",
            "file_url": None,
            "session_id": session_id or str(uuid.uuid4()),
        })


@app.route("/api/feedback", methods=["POST"])
def api_feedback():
    """Optional thumbs / note for self-iteration signals (stored in SQLite)."""
    data = request.json or {}
    session_id = data.get("session_id") or ""
    rating = data.get("rating")
    note = data.get("note") or ""
    if rating is not None:
        try:
            rating = int(rating)
        except (TypeError, ValueError):
            rating = 0
    memory_store.record_feedback(session_id, rating, note)
    return jsonify({"ok": True})


@app.route("/api/memory/summary")
def api_memory_summary():
    token = os.environ.get("MEMORY_ADMIN_TOKEN")
    if token and request.headers.get("X-Memory-Token") != token:
        return jsonify({"error": "unauthorized"}), 401
    return jsonify(memory_store.memory_summary_json())


@app.route('/api/batch', methods=['POST'])
def batch_api():
    """Upload a PPTX → LLM re-interprets each slide → rebuild all into one branded PPTX."""
    try:
        uploaded = request.files.get('file')
        if not uploaded or not uploaded.filename.lower().endswith('.pptx'):
            return jsonify({"reply": "Please upload a .pptx file.", "file_url": None})

        upload_dir = os.path.join(os.path.dirname(__file__), 'output', 'uploads')
        os.makedirs(upload_dir, exist_ok=True)
        tmp_name = f"{uuid.uuid4().hex[:8]}_{uploaded.filename}"
        tmp_path = os.path.join(upload_dir, tmp_name)
        uploaded.save(tmp_path)
        app.logger.info(f"[batch] file={uploaded.filename}")

        src = ReadPptx(tmp_path)
        slide_texts = []
        for i, slide in enumerate(src.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        t = p.text.strip()
                        if t:
                            texts.append(t)
            slide_texts.append(f"[Slide {i+1}]\n" + "\n".join(texts))

        all_content = "\n\n".join(slide_texts)

        prompt = (
            f"Below is a {len(src.slides)}-slide presentation. "
            f"For EACH slide, output a JSON object choosing the best layout_type. "
            f"Return a JSON ARRAY of objects, one per slide. "
            f"Preserve ALL original text — do not summarize, shorten or omit any content. "
            f"Keep the exact wording from the source.\n\n"
            f"--- Source slides ---\n{all_content[:12000]}\n--- End ---"
        )

        response = messages_create_with_fallback(
            max_tokens=8192,
            system=SYSTEM_PROMPT.replace(
                "Output exactly ONE JSON object",
                "Output a JSON ARRAY of objects, one per slide"
            ),
            messages=[{"role": "user", "content": prompt}],
        )

        raw = response.content[0].text.strip()
        cleaned = raw
        if cleaned.startswith("```"):
            cleaned = cleaned.split("\n", 1)[1]
            if cleaned.endswith("```"):
                cleaned = cleaned[:-3]
        parsed = json.loads(cleaned)

        if isinstance(parsed, dict):
            slides_data = parsed.get("slides", [parsed])
        elif isinstance(parsed, list):
            slides_data = parsed
        else:
            raise ValueError("Expected JSON array or object with 'slides' key.")

        from builder import build_cards, build_cover, build_big_numbers, build_case_study, SLIDE_W, SLIDE_H
        from pptx import Presentation as NewPrs

        prs = NewPrs()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        builders = {
            "cards": build_cards,
            "big_numbers": build_big_numbers,
            "cover": build_cover,
            "case_study": build_case_study,
        }

        for sd in slides_data:
            if not isinstance(sd, dict):
                continue
            lt = sd.get("layout_type", "cards")
            fn = builders.get(lt, build_cards)
            fn(prs, sd)

        out_name = f"batch-{uuid.uuid4().hex[:8]}.pptx"
        out_path = os.path.join(os.path.dirname(__file__), 'output', out_name)
        os.makedirs(os.path.dirname(out_path), exist_ok=True)
        prs.save(out_path)

        app.logger.info(f"[batch] OK slides={len(slides_data)} file={out_name}")
        return jsonify({
            "reply": f"Batch complete! {len(slides_data)} slides generated in Dify brand style.",
            "file_url": f"/download/{out_name}",
        })

    except Exception as e:
        app.logger.error(f"[batch] ERROR: {e}", exc_info=True)
        return jsonify({
            "reply": f"Batch error: {e}",
            "file_url": None,
        })


@app.route('/download/<filename>')
def download(filename):
    path = os.path.join(os.path.dirname(__file__), 'output', filename)
    if os.path.exists(path):
        return send_file(path, as_attachment=True)
    return "File not found", 404


if __name__ == '__main__':
    os.makedirs('output', exist_ok=True)
    port = int(os.environ.get('PORT', 8080))
    print(f"\n  Dify PPT Generator running at: http://localhost:{port}\n")
    app.run(host='0.0.0.0', port=port, debug=False)
